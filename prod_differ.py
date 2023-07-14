
# Online Python - IDE, Editor, Compiler, Interpreter

import pandas as pd
import sqlite3
import streamlit as st
from bokeh.models import DatetimeTickFormatter, HoverTool,Title,LinearAxis, Range1d,NumeralTickFormatter,ColumnDataSource,LabelSet
from bokeh.plotting import figure
import datetime
import math
import pandas as pd
from bokeh.palettes import Category10
import numpy as np
import plotly.graph_objects as go
from pptx import Presentation
from pptx.util import Inches
import io
import plotly.io as pio
import os


def differ(): 
    option1 = st.checkbox('actual')
    option2 = st.checkbox('well test')
    

# Check if Option 1 is selected
    if option1:
       dh = pd.read_csv('differ.csv',low_memory=False)
       extracted_string = 'actual differ'
    elif option2:
       dh = pd.read_csv('TestDiffer.csv',low_memory=False)
       extracted_string = 'test'
    # Check if Option 2 is selected
    else:
       dh = pd.read_csv('differ.csv',low_memory=False)
       extracted_string = 'actual differ' 
    df =dh               
          
  
    
    df['Datee'] = pd.to_datetime(df['Datee'], errors='coerce')
    col1, col2,col3,col4,col5,col6,col7 = st.columns([1,1,1,1,1,1,1])
    unique_dates = df['Datee'].dt.strftime('%Y-%m-%d').unique().tolist()  # Convert dates to desired format
    
    
    def merge(df1,df2):
        merged_df = pd.merge(df1, df2, left_on='WellName', right_on='WellName', how='outer')
        merged_df['net_differ'] = merged_df['NetOil_x'].fillna(0) - merged_df['NetOil_y'].fillna(0)
        merged_df['gross_differ'] = merged_df['Gross_x'].fillna(0) - merged_df['Gross_y'].fillna(0)
        
               

        return merged_df
        
        
    def differ_field(merged_df):
        df_differ = merged_df.iloc[:,[3,4,10,11,7,15,16,14]]
        df_filtered = df_differ.groupby('Field_x')[['Gross_x','NetOil_x','net_differ','gross_differ']].sum().reset_index()
        df_filtered2 = df_differ.groupby('Field_y')[['Gross_y','NetOil_y']].sum().reset_index()
        df_final = pd.merge(df_filtered, df_filtered2, left_on='Field_x', right_on='Field_y', how='left')
        df_filtered = df_final.iloc[:,[0,1,2,6,7,3,4]]
        columns_to_convert = ['Gross_x', 'NetOil_x', 'Gross_y', 'NetOil_y', 'net_differ', 'gross_differ']
        df_filtered[columns_to_convert] = df_filtered[columns_to_convert].astype(float)
        # Calculate the total sum for all columns starting from "Gross_x"
        total_sum = df_filtered.loc[:, "Gross_x":].sum(axis=0)
        # Add the total sum as a new row to the dataframe
        total_row = {"Field_x": "Total"}
        total_row.update(total_sum)
        total_row_df = pd.DataFrame(total_row, index=[0])
        df_filtered = pd.concat([df_filtered, total_row_df], ignore_index=True)
    
        return df_filtered
    
    
    
    def select_differ(merged_df,field):
        df_field = merged_df[(merged_df['Field_x'] == field) & (abs(merged_df['net_differ']) >35)]
        df_field = df_field.iloc[:,[2,15,16]]
        return df_field
    
    
    
    
    def create_column_chart(df,title):
    # Create the figure
        fig = go.Figure()
    
        # Add the net difference bar
        fig.add_trace(go.Bar(
            x=df['WellName'],
            y=df['net_differ'],
            name='Net Difference',
            marker_color='green'
        ))
    
        # Add the gross difference bar
        fig.add_trace(go.Bar(
            x=df['WellName'],
            y=df['gross_differ'],
            name='Gross Difference',
            marker_color='orange'
        ))
    
        # Add data labels
        for x, y in zip(df['WellName'], df['net_differ']):
            fig.add_annotation(
                x=x,
                y=y,
                text="<b>{}</b>".format(int(y)),
                showarrow=False,
                font=dict(size=12,color='black'),
                xanchor='center',
                yanchor='bottom'
            )
    
        # Customize the layout
        fig.update_layout(
       title=title,
       xaxis_title='Well Name',
       yaxis_title='Difference',
       barmode='group',
       width=500  # Adjust the width here
   )
    
        # Show the chart in Streamlit
        return fig
        
        
    def create_powerpoint_presentation(figures, df, title, subtitle):
    # Create the PowerPoint presentation
        prs = Presentation()

        # Add the first slide with the specified title and subtitle
        slide_layout = prs.slide_layouts[0]  # Use the layout for the title slide
        slide = prs.slides.add_slide(slide_layout)
        title_placeholder = slide.shapes.title
        title_placeholder.text = str(title)
        subtitle_placeholder = slide.placeholders[1]
        subtitle_placeholder.text = str(subtitle)

        # Add a slide for the dataframe as a table
        slide_layout = prs.slide_layouts[6]  # Use the layout for a blank slide
        slide = prs.slides.add_slide(slide_layout)

        # Add the dataframe to the slide as a table
        left = top = Inches(1)  # Adjust the position of the table on the slide
        width = Inches(8)
        height = Inches(6)  # Adjust the size of the table
        slide.shapes.add_table(rows=len(df) + 1, cols=len(df.columns), left=left, top=top, width=width, height=height)

        # Populate the table with the dataframe data
        table = slide.shapes[-1].table
        # Set the column headers
        for i, column in enumerate(df.columns):
            table.cell(0, i).text = column
        # Set the cell values
        for i, row in enumerate(df.itertuples(index=False)):
            for j, value in enumerate(row):
                table.cell(i + 1, j).text = str(value)

        for fig in figures:
            # Save the chart as an image
            chart_image_path = 'chart.png'
            fig.write_image(chart_image_path)

            # Add a slide to the presentation
            slide_layout = prs.slide_layouts[6]  # Use the layout for a blank slide
            slide = prs.slides.add_slide(slide_layout)

            # Add the chart image to the slide
            left = top = Inches(0)
            slide.shapes.add_picture(chart_image_path, left, top, prs.slide_width, prs.slide_height)

        # Save the PowerPoint presentation
        presentation_path = 'presentation.pptx'
        prs.save(presentation_path)
        os.startfile(presentation_path)

        print("Presentation created successfully: {}".format(presentation_path))

   
   
    with col1:
        default1 = ["2023-06-19"]
        default2 = ["2023-03-26"]
        selected_date1 = st.multiselect("Select a date (1)", unique_dates,default1, key="date1")
        checkbox_value = st.checkbox("Field differ",value=True)
       
        
    with col2:
        selected_date2 = st.multiselect("Select a date (2)", unique_dates,default2,key="date2")
        
        

    with col3:
        checkbox_value2 = st.checkbox("Abrar differ",value=True)
        
        
    with col4:
        checkbox_value3 = st.checkbox("AbrarSouth differ",value=True)    
        
    with col5:
        checkbox_value4 = st.checkbox("Ganna differ",value=True)    
        
        
    with col6:
        checkbox_value5 = st.checkbox("Sidra differ",value=True)    
        
        
    with col7:
        export_ppt=st.button('Export ppt')
        show_well = st.button('show well differ')
        


    
           
      
        
    
    
    if selected_date1 and selected_date2:
        df['RunTime'] = pd.to_numeric(df['RunTime'], errors='coerce')
        df["RunTime"] = df["RunTime"].astype(float)
        df = df[df['RunTime']>0]
        df1 = df[df['Datee'].isin(selected_date1)]
        df2 = df[df['Datee'].isin(selected_date2)]
        merged_df=merge(df1,df2)
        df_filtered = differ_field(merged_df)
        df_abrar = select_differ(merged_df, 'Abrar')
        df_as = select_differ(merged_df, 'Abrar south')
        df_ganna = select_differ(merged_df, 'Ganna')
        df_sidra = select_differ(merged_df, 'Sidra')
        
        
    else: 
        df1 = None
        df2 = None

        
        
    if checkbox_value:
        if df1 is not None and df2 is not None:
            with col1:         
              st.dataframe(df_filtered)
    else:
        st.write(" ")
        
        
        
    if checkbox_value2:
        if df1 is not None and df2 is not None:
            with col2:         
              a=create_column_chart(df_abrar,'Abrar')
              st.plotly_chart(a)
    else:
        st.write(" ")
        
        
        
    if checkbox_value3:
        if df1 is not None and df2 is not None:
                      
            with col2:
               b= create_column_chart(df_as,'Abrar_south')
               st.plotly_chart(b)
    else:
        st.write(" ")    
        
    
    
    
    if checkbox_value4:
        if df1 is not None and df2 is not None:
                      
            with col5:
                for i in range(3):
                    st.write(" ")
                c=create_column_chart(df_ganna,'Ganna')
                st.plotly_chart(c)
    else:
        st.write(" ")    
       
        
       
        
    if checkbox_value5:
        if df1 is not None and df2 is not None:
                      
            with col5:
                for i in range(1):
                    st.write(" ")
                d=create_column_chart(df_sidra,'Sidra')
                st.plotly_chart(d)
    else:
        st.write(" ")    
      
   
    
    if show_well:
        if df1 is not None and df2 is not None:
                      
            with col7:
               numeric_columns = merged_df.select_dtypes(include=np.number).columns
               merged_df[numeric_columns] = merged_df[numeric_columns].fillna(0).astype(int)

               merged_df = merged_df.sort_values('net_differ') 
               st.dataframe( merged_df)
               
    else:
        st.write(" ")    


    if export_ppt:
        if df1 is not None and df2 is not None:
                      
            with col7:
                a=create_column_chart(df_abrar,'Abrar')
                b= create_column_chart(df_as,'Abrar_south')
                c=create_column_chart(df_ganna,'Ganna')
                d=create_column_chart(df_sidra,'Sidra')
                


                figures = [a,b,c,d]  # List of Plotly figures
                name = str(selected_date1[0])+" " + "vs" + " "+ str(selected_date2[0])
                sub_name = extracted_string + "production differ"
                create_powerpoint_presentation(figures,df_filtered,name,sub_name)
                            
    else:
       st.write(" ")    
    


    

 
    